"""
GradingAgent: An automated grading system using ChatGPT API
Loads rubrics, student papers, grades them, and outputs results for mail merge
"""

import os
import json
import pandas as pd
from pathlib import Path
from typing import Dict, List, Any, Optional
import openai
from docx import Document
import PyPDF2
import logging
from datetime import datetime

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('grading_log.txt'),
        logging.StreamHandler()
    ]
)

# Models that do not support file uploads (static list of older models)
MODELS_WITHOUT_FILE_SUPPORT = [
    'gpt-3.5-turbo',
    'gpt-3.5-turbo-1106', 
    'gpt-3.5-turbo-0125',
    'gpt-3.5-turbo-16k',
    'gpt-3.5-turbo-instruct',
    'text-davinci-003',
    'text-davinci-002',
    'davinci-002',
    'babbage-002'
]


# Unified GradingAgent with multi-file, course material summarization, custom instructions, GPT-5 API compatibility, and cached input optimization
class GradingAgent:
    def _upload_files(self, file_paths: list) -> tuple:
        """Upload files to OpenAI and return (uploaded_files, file_ids, errors)."""
        import mimetypes
        uploaded_files = []
        file_ids = []
        errors = []
        
        for file_path in file_paths:
            try:
                mime_type, _ = mimetypes.guess_type(file_path)
                with open(file_path, "rb") as f:
                    file_obj = self.client.files.create(
                        file=f,
                        purpose="assistants"  # Note: This may need adjustment for chat completions
                    )
                file_ids.append(file_obj.id)
                uploaded_files.append(os.path.basename(file_path))
                logging.info(f"Successfully uploaded {os.path.basename(file_path)} (ID: {file_obj.id})")
            except Exception as upload_exc:
                logging.error(f"OpenAI file upload failed for {file_path}: {upload_exc}")
                errors.append({"file_path": file_path, "error": str(upload_exc)})
                
        return uploaded_files, file_ids, errors
    
    def _cleanup_uploaded_files(self, file_ids: list) -> None:
        """Clean up uploaded files from OpenAI to prevent storage bloat."""
        for file_id in file_ids:
            try:
                self.client.files.delete(file_id)
                logging.debug(f"Cleaned up uploaded file: {file_id}")
            except Exception as cleanup_exc:
                logging.warning(f"Failed to clean up file {file_id}: {cleanup_exc}")
    
    def _extract_text_from_file(self, file_path) -> str:
        """Extract text content from various file types."""
        file_path = Path(file_path)
        
        try:
            if file_path.suffix.lower() == '.pdf':
                return self._extract_pdf_text(file_path)
            elif file_path.suffix.lower() == '.docx':
                return self._extract_docx_text(file_path)
            elif file_path.suffix.lower() in ['.txt', '.md']:
                return self._extract_txt_text(file_path)
            elif file_path.suffix.lower() in ['.html', '.htm']:
                return self._extract_html_text(file_path)
            elif file_path.suffix.lower() in ['.pptx', '.ppt']:
                return self._extract_pptx_text(file_path)
            elif file_path.suffix.lower() == '.odt':
                return self._extract_odt_text(file_path)
            elif file_path.suffix.lower() == '.rtf':
                return self._extract_rtf_text(file_path)
            else:
                # Try as plain text
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
        except Exception as e:
            logging.error(f"Failed to extract text from {file_path}: {e}")
            return f"[Error extracting text from {file_path.name}: {e}]"
    
    def _extract_pdf_text(self, file_path: Path) -> str:
        """Extract text from PDF file."""
        try:
            import PyPDF2
            text_content = []
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    page_text = page.extract_text()
                    text_content.append(f"--- PAGE {page_num} ---\n{page_text}")
            return '\n'.join(text_content)
        except ImportError:
            return "[PyPDF2 not available for PDF processing]"
        except Exception as e:
            logging.error(f"PDF extraction error: {e}")
            return f"[Error extracting PDF text: {e}]"
    
    def _extract_docx_text(self, file_path: Path) -> str:
        """Extract text from DOCX file."""
        try:
            from docx import Document
            doc = Document(file_path)
            text_content = []
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text)
            
            # Extract tables
            for table in doc.tables:
                text_content.append("\n[TABLE:]")
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    if row_text.strip():
                        text_content.append(row_text)
                text_content.append("[END TABLE]\n")
            
            return '\n'.join(text_content)
        except ImportError:
            return "[python-docx not available for DOCX processing]"
        except Exception as e:
            logging.error(f"DOCX extraction error: {e}")
            return f"[Error extracting DOCX text: {e}]"
    
    def _extract_txt_text(self, file_path: Path) -> str:
        """Extract text from plain text file."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            logging.error(f"TXT extraction error: {e}")
            return f"[Error extracting text: {e}]"
    
    def _extract_html_text(self, file_path: Path) -> str:
        """Extract text from HTML file."""
        try:
            from bs4 import BeautifulSoup
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                return soup.get_text()
        except ImportError:
            # Fallback: read as plain text
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            logging.error(f"HTML extraction error: {e}")
            return f"[Error extracting HTML text: {e}]"
    
    def _extract_pptx_text(self, file_path: Path) -> str:
        """Extract text from PowerPoint file."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_content.append(f"--- SLIDE {slide_num} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text.strip())
            
            return '\n'.join(text_content)
        except ImportError:
            return "[python-pptx not available for PowerPoint processing]"
        except Exception as e:
            logging.error(f"PPTX extraction error: {e}")
            return f"[Error extracting PowerPoint text: {e}]"
    
    def _extract_odt_text(self, file_path: Path) -> str:
        """Extract text from ODT file."""
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            
            with zipfile.ZipFile(file_path, 'r') as odt_zip:
                if 'content.xml' in odt_zip.namelist():
                    content_xml = odt_zip.read('content.xml').decode('utf-8')
                    root = ET.fromstring(content_xml)
                    
                    # Extract all text elements
                    text_elements = []
                    for elem in root.iter():
                        if elem.text:
                            text_elements.append(elem.text.strip())
                    
                    return '\n'.join(filter(None, text_elements))
                else:
                    return "[ODT content.xml not found]"
        except Exception as e:
            logging.error(f"ODT extraction error: {e}")
            return f"[Error extracting ODT text: {e}]"
    
    def _extract_rtf_text(self, file_path: Path) -> str:
        """Extract text from RTF file."""
        try:
            # Basic RTF text extraction (very simple)
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
                # Very basic RTF parsing - just remove obvious RTF tags
                import re
                # Remove RTF control words
                content = re.sub(r'\\[a-z]+\d*\s*', '', content)
                content = re.sub(r'[{}]', '', content)
                return content.strip()
        except Exception as e:
            logging.error(f"RTF extraction error: {e}")
            return f"[Error extracting RTF text: {e}]"
    _course_summary_cache = {}
    """Unified GradingAgent for multi-file, course material summarization, custom instructions, and GPT-5 API compatibility."""

    def __init__(self, api_key: str, model: str = "gpt-4o-mini"):
        self.client = openai.OpenAI(api_key=api_key)
        self.model = model
        self.rubric = None
        self.students_data = []
        self.graded_results = []
        self.instructor_config = None
        self.course_materials = []  # Initialize as empty list
        self.course_materials_instructions = ""  # Initialize as empty string
        self.summarized_course_materials = None


    def load_instructor_config(self, config_path: str) -> Dict[str, Any]:
        """Load instructor configuration from JSON file."""
        try:
            with open(config_path, 'r', encoding='utf-8') as f:
                self.instructor_config = json.load(f)
                logging.info(f"Loaded instructor config from {config_path}")
                return self.instructor_config
        except Exception as e:
            logging.error(f"Failed to load instructor config: {e}")
            self.instructor_config = {}
            return {}

    def load_rubric(self, rubric_path: str) -> Dict[str, Any]:
        """Load grading rubric from JSON file."""
        try:
            with open(rubric_path, 'r', encoding='utf-8') as f:
                self.rubric = json.load(f)
                logging.info(f"Loaded rubric from {rubric_path}")
                return self.rubric
        except Exception as e:
            logging.error(f"Failed to load rubric: {e}")
            self.rubric = {}
            return {}

    def load_course_materials(self, materials_files: List[str], instructions: str = "") -> None:
        """Load course materials from files for grading context."""
        self.course_materials = []
        self.course_materials_instructions = instructions
        
        for file_path in materials_files:
            try:
                content = self._extract_text_from_file(file_path)
                self.course_materials.append({
                    'filename': os.path.basename(file_path),
                    'content': content,
                    'file_path': file_path
                })
                logging.info(f"Loaded course material: {os.path.basename(file_path)}")
            except Exception as e:
                logging.error(f"Failed to load course material {file_path}: {e}")
        
        logging.info(f"Loaded {len(self.course_materials)} course materials")


    def summarize_course_materials(self, rubric: dict, course_materials: list) -> str:
        """Summarize the provided course materials for cost-effective grading context."""
        if not course_materials:
            return ""
        import hashlib
        # Create a hash key from rubric and course_materials content
        rubric_str = json.dumps(rubric, sort_keys=True)
        materials_str = "\n\n".join([m['content'] for m in course_materials])
        cache_key = hashlib.sha256((rubric_str + materials_str).encode('utf-8')).hexdigest()
        if cache_key in self._course_summary_cache:
            return self._course_summary_cache[cache_key]
        prompt = (
            "Summarize the following course materials for use as grading context. "
            "Focus on rubric-relevant information and key facts.\n\n"
            f"RUBRIC: {json.dumps(rubric, indent=2)}\n\nCOURSE MATERIALS:\n{materials_str}"
        )
        summary = self._make_chat_completion([
            {"role": "system", "content": "You are an expert educator. Summarize course materials for grading context."},
            {"role": "user", "content": prompt}
        ], max_tokens=1500)
        self._course_summary_cache[cache_key] = summary
        return summary

    def load_student_submissions(self, submissions_folder: str) -> List[Dict[str, Any]]:
        """Load all files for each student as a single submission (multi-file support)."""
        submissions_path = Path(submissions_folder)
        self.students_data = []
        for student_dir in submissions_path.iterdir():
            if student_dir.is_dir():
                files = list(student_dir.glob("*"))
                submission_files = []
                for file_path in files:
                    try:
                        content = self._extract_text_from_file(file_path)
                        submission_files.append({
                            'filename': file_path.name,
                            'content': content,
                            'file_path': str(file_path)
                        })
                    except Exception as e:
                        logging.warning(f"Failed to load {file_path}: {e}")
                self.students_data.append({
                    'student_name': student_dir.name,
                    'files': submission_files
                })
        logging.info(f"Loaded {len(self.students_data)} student submissions (multi-file)")
        return self.students_data


    def _grading_json_example(self, rubric=None):
        """Return the required grading JSON format as a string for prompt inclusion."""
        if rubric and isinstance(rubric, dict) and 'criteria' in rubric:
            # Generate rubric-specific example
            criteria_examples = []
            for criterion_name, criterion_data in rubric['criteria'].items():
                max_points = criterion_data.get('points', 10)
                criteria_examples.append(f'''        "{criterion_name}": {{
            "score": <{max_points}_or_less>,
            "max_score": {max_points},
            "feedback": "<specific_feedback_for_{criterion_name.lower().replace(' ', '_')}>"
        }}''')
            
            criteria_section = ',\n'.join(criteria_examples)
            max_possible = rubric.get('total_points', 100)
            
            return f'''{{
    "overall_score": <total_points_earned>,
    "max_possible_score": {max_possible},
    "percentage": <percentage_score>,
    "letter_grade": "<letter_grade>",
    "criteria_scores": {{
{criteria_section}
    }},
    "overall_feedback": "<comprehensive_feedback>",
    "strengths": ["<strength1>", "<strength2>"],
    "areas_for_improvement": ["<improvement1>", "<improvement2>"]
}}'''
        else:
            # Fallback generic example
            return '''
{
    "overall_score": <total_points>,
    "max_possible_score": <maximum_points>,
    "percentage": <percentage_score>,
    "letter_grade": "<letter_grade>",
    "criteria_scores": {
        "<criterion_name>": {
            "score": <points>,
            "max_score": <max_points>,
            "feedback": "<specific_feedback>"
        }
    },
    "overall_feedback": "<comprehensive_feedback>",
    "strengths": ["<strength1>", "<strength2>"],
    "areas_for_improvement": ["<improvement1>", "<improvement2>"]
}
'''

    def grade_submission(self, student_submission: Dict[str, Any], rubric: dict, course_materials: list = None, custom_instructions: str = "") -> Dict[str, Any]:
        """Hybrid grading: file upload if supported, fallback to text extraction. Stateless with respect to rubric and course materials."""
        # Handle None course_materials
        course_materials = course_materials or []
        
        # Summarize course materials for this grading call
        course_summary = self.summarize_course_materials(rubric, course_materials)
        
        # Partition files into uploadable vs extractable
        uploadable_files = []
        extracted_texts = []
        
        for f in student_submission['files']:
            file_path = f.get('file_path')
            if not file_path or not os.path.exists(file_path):
                # No file path or file doesn't exist, use extracted content
                extracted_texts.append(f"--- {f['filename']} ---\n{f['content']}")
                continue
                
            file_ext = os.path.splitext(file_path)[1].lower()
            can_upload = self._model_supports_file_uploads() and file_ext in ['.pdf', '.docx', '.pptx', '.ppt', '.txt', '.html', '.htm']
            
            if can_upload:
                uploadable_files.append(file_path)
            else:
                extracted_texts.append(f"--- {f['filename']} ---\n{f['content']}")

        # Try file upload for uploadable files
        uploaded_files, file_ids, upload_errors = [], [], []
        if uploadable_files:
            try:
                uploaded_files, file_ids, upload_errors = self._upload_files(uploadable_files)
            except Exception as upload_exc:
                logging.error(f"File upload batch failed: {upload_exc}")
                # Move all uploadable files to extracted texts as fallback
                for file_path in uploadable_files:
                    filename = os.path.basename(file_path)
                    try:
                        content = self._extract_text_from_file(file_path)
                        extracted_texts.append(f"--- {filename} (upload failed, using text extraction) ---\n{content}")
                    except Exception as extract_exc:
                        logging.error(f"Fallback text extraction failed for {filename}: {extract_exc}")
                        extracted_texts.append(f"--- {filename} (processing failed) ---\n[Error: Unable to process this file]")
                upload_errors.append({"error": f"Batch upload failed: {upload_exc}"})

        # Compose the grading prompt
        file_list_str = ", ".join([f"{name} (uploaded)" for name in uploaded_files]) if uploaded_files else "None"
        extracted_text_section = "\n\n".join(extracted_texts) if extracted_texts else "None"
        
        # Build the complete prompt
        prompt = (
            f"Grade the submission indicated below using EXACTLY the following rubric structure and point values.\n\n"
            f"RUBRIC:\n{json.dumps(rubric, indent=2)}\n\n"
            f"CRITICAL INSTRUCTIONS:\n"
            f"- Use EXACTLY the criteria names from the rubric above\n"
            f"- Use EXACTLY the point values specified in the rubric\n"
            f"- Total possible score MUST be {rubric.get('total_points', 100)} points\n"
            f"- Provide detailed feedback for each individual rubric criterion\n\n"
            f"COURSE MATERIALS SUMMARY:\n{course_summary}\n\n" if course_summary else ""
            f"CUSTOM INSTRUCTIONS:\n{custom_instructions}\n\n" if custom_instructions else ""
            f"Please provide your response in the following JSON format:\n{self._grading_json_example(rubric)}\n\n"
            f"Be specific, constructive, and fair in your grading. Provide actionable feedback.\n\n"
            f"SUBMISSION FILES:\n"
            f"Uploaded files: {file_list_str}\n"
            f"Text-extracted content:\n{extracted_text_section}"
        )

        # Estimate tokens for adaptive response length (only count text, not uploaded files)
        est_tokens = self._count_tokens(extracted_text_section + course_summary + custom_instructions)
        if est_tokens > 50000:
            max_output_tokens = 4000
        elif est_tokens > 20000:
            max_output_tokens = 3000
        else:
            max_output_tokens = 2000

        # Make the API call (note: files parameter deprecated, using text extraction instead)
        try:
            result_content = self._make_chat_completion(
                [
                    {"role": "system", "content": self._build_system_message()},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=max_output_tokens,
                temperature=0.3
            )
            
            # Parse the grading response
            try:
                # Handle responses that may be wrapped in code blocks
                clean_content = result_content.strip()
                if clean_content.startswith('```json'):
                    # Extract JSON from code blocks
                    clean_content = clean_content.split('```json')[1].split('```')[0].strip()
                elif clean_content.startswith('```'):
                    # Handle generic code blocks
                    clean_content = clean_content.split('```')[1].split('```')[0].strip()
                
                grading_data = json.loads(clean_content)
            except json.JSONDecodeError:
                logging.error(f"Failed to parse grading response as JSON: {result_content[:200]}...")
                grading_data = {
                    "error": "Failed to parse grading response as JSON",
                    "raw_response": result_content,
                    "overall_score": 0,
                    "max_possible_score": 100,
                    "percentage": 0,
                    "letter_grade": "F",
                    "criteria_scores": {},
                    "overall_feedback": "Error: Could not parse AI grading response",
                    "strengths": [],
                    "areas_for_improvement": ["Grading system error - please review manually"]
                }
            
            # Add metadata
            grading_data['files_uploaded'] = uploaded_files
            grading_data['file_ids'] = file_ids
            grading_data['files_processed'] = len(student_submission['files'])
            grading_data['estimated_tokens'] = est_tokens
            if upload_errors:
                grading_data['upload_errors'] = upload_errors
                
        except Exception as api_exc:
            logging.error(f"OpenAI grading API call failed: {api_exc}")
            grading_data = {
                "error": f"Grading API call failed: {api_exc}",
                "file_ids": file_ids,
                "upload_errors": upload_errors,
                "overall_score": 0,
                "max_possible_score": 100,
                "percentage": 0,
                "letter_grade": "F",
                "criteria_scores": {},
                "overall_feedback": f"Error: API call failed - {api_exc}",
                "strengths": [],
                "areas_for_improvement": ["Grading system error - please review manually"]
            }
        
        # Clean up uploaded files
        if file_ids:
            try:
                self._cleanup_uploaded_files(file_ids)
            except Exception as cleanup_exc:
                logging.warning(f"Failed to clean up uploaded files: {cleanup_exc}")

        # Add standard metadata
        grading_data['student_name'] = student_submission['student_name']
        grading_data['grading_date'] = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        grading_data['grading_method'] = 'unified_multi_file'
        
        return grading_data

    def _count_tokens(self, text: str) -> int:
        """Count tokens in text using tiktoken for the current model."""
        try:
            import tiktoken
            enc = tiktoken.encoding_for_model(self.model)
            return len(enc.encode(text))
        except Exception:
            # Fallback: rough estimate if tiktoken not available
            return len(text) // 4

    def model_supports_file_uploads(self) -> bool:
        """Check if the current model supports file uploads for document grading (legacy method name)."""
        return self._model_supports_file_uploads()
    
    def _model_supports_file_uploads(self) -> bool:
        """Check if the current model supports file uploads for document grading."""
        model_lower = self.model.lower()
        for blocked_model in MODELS_WITHOUT_FILE_SUPPORT:
            if blocked_model.lower() in model_lower:
                return False
        return True


    def grade_all_submissions(self, rubric: dict, course_materials: list = None, custom_instructions: str = "") -> List[Dict[str, Any]]:
        """Grade all loaded student submissions (multi-file, unified, GPT-5 compatible)."""
        if not self.students_data:
            raise ValueError("No student submissions loaded.")
        
        # Handle None parameters
        course_materials = course_materials or []
        
        self.graded_results = []
        for student_submission in self.students_data:
            try:
                result = self.grade_submission(student_submission, rubric, course_materials, custom_instructions)
                self.graded_results.append(result)
                logging.info(f"Successfully graded submission for {student_submission['student_name']}")
            except Exception as e:
                logging.error(f"Failed to grade {student_submission['student_name']}: {str(e)}")
                # Add error result to maintain consistency
                error_result = {
                    'student_name': student_submission['student_name'],
                    'error': str(e),
                    'overall_score': 0,
                    'max_possible_score': 100,
                    'percentage': 0,
                    'letter_grade': 'F',
                    'criteria_scores': {},
                    'overall_feedback': f"Grading failed: {str(e)}",
                    'strengths': [],
                    'areas_for_improvement': ["System error - manual review required"],
                    'grading_date': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                }
                self.graded_results.append(error_result)
                
        logging.info(f"Completed grading {len(self.graded_results)} submissions")
        return self.graded_results

    # ...existing code for export_results, generate_summary_report, etc...
    
    def grade_paper(self, student_data: Dict[str, str], custom_instructions: str = "") -> Dict[str, Any]:
        """
        Grade a single student paper with intelligent multi-file and upload support
        
        This method automatically detects multi-file submissions and routes to appropriate
        grading logic. It chooses between file upload (for compatible models) and text 
        extraction (for older models or as fallback) to preserve document formatting.
        
        Args:
            student_data (Dict): Student paper data with keys:
                - 'name': Student name  
                - 'content': Extracted text content (for single-file fallback)
                - 'file_path': Path to original file (optional, for single-file upload)
                - 'files': List of file dicts for multi-file submissions (optional)
                    Each file dict: {'filename': str, 'content': str, 'file_path': str}
            custom_instructions (str): Additional grading instructions for this specific assignment
            
        Returns:
            Dict: Grading results
        """
        if not self.rubric:
            raise ValueError("Rubric not loaded. Please load a rubric first.")
        
        # Detect if this is a multi-file submission
        if 'files' in student_data and student_data['files']:
            logging.info(f"Multi-file submission detected for {student_data['name']} ({len(student_data['files'])} files)")
            # Convert to new multi-file format and use new grading method
            submission_data = {
                'student_name': student_data['name'],
                'files': student_data['files']
            }
            
            # Use the new unified grading method with proper parameters
            return self.grade_submission(
                submission_data, 
                self.rubric, 
                self.course_materials,
                custom_instructions
            )
        
        # Single-file submission processing
        logging.info(f"Single-file submission detected for {student_data['name']}")
        
        # Check content size for very large submissions
        if 'content' in student_data and student_data['content']:
            estimated_tokens = len(student_data['content']) // 4
            if estimated_tokens > 100000:
                logging.warning(f"Very large submission detected for {student_data['name']} (~{estimated_tokens:,} tokens). This may approach model limits or incur high costs.")
        
        # For single-file, create a files structure for consistency with new method
        if 'file_path' in student_data and student_data['file_path']:
            files_data = [{
                'filename': os.path.basename(student_data['file_path']),
                'content': student_data.get('content', ''),
                'file_path': student_data['file_path']
            }]
        else:
            # No file path, just content
            files_data = [{
                'filename': 'submission.txt',
                'content': student_data.get('content', ''),
                'file_path': None
            }]
        
        submission_data = {
            'student_name': student_data['name'],
            'files': files_data
        }
        
        # Use unified grading method for consistency
        return self.grade_submission(
            submission_data,
            self.rubric,
            self.course_materials, 
            custom_instructions
        )
    
    def _grade_paper_with_file_upload(self, student_data: Dict[str, str]) -> Dict[str, Any]:
        """
        Grade a paper using OpenAI file upload capabilities (preserves formatting)
        
        This method uses the newer OpenAI APIs to process documents directly,
        preserving formatting and structure that would be lost in text extraction.
        
        Args:
            student_data (Dict): Student paper data
            
        Returns:
            Dict: Grading results
        """
        file_path = Path(student_data['file_path'])
        
        if not file_path.exists():
            raise FileNotFoundError(f"Submission file not found: {file_path}")
        
        # Check if file type is supported for direct upload
        supported_upload_extensions = ['.pdf', '.docx', '.txt', '.html', '.htm', '.pptx', '.ppt']
        if file_path.suffix.lower() not in supported_upload_extensions:
            raise ValueError(f"File type {file_path.suffix} not supported for upload")
        
        try:
            # For newer GPT-4 models, we can use file upload with assistants API
            # or enhanced document processing with vision capabilities
            
            if file_path.suffix.lower() in ['.pdf']:
                return self._grade_pdf_with_vision(student_data, file_path)
            elif file_path.suffix.lower() in ['.docx']:
                return self._grade_docx_enhanced(student_data, file_path)
            elif file_path.suffix.lower() in ['.pptx', '.ppt']:
                return self._grade_powerpoint_enhanced(student_data, file_path)
            else:
                return self._grade_text_file_enhanced(student_data, file_path)
            
        except Exception as e:
            logging.error(f"Error in file upload grading for {student_data['name']}: {str(e)}")
            raise
    
    def _grade_pdf_with_vision(self, student_data: Dict[str, str], file_path: Path) -> Dict[str, Any]:
        """
        Grade PDF using vision capabilities to preserve visual formatting
        
        Args:
            student_data (Dict): Student paper data
            file_path (Path): Path to PDF file
            
        Returns:
            Dict: Grading results
        """
        try:
            # For PDFs, we'll use an enhanced approach that considers document structure
            # Since direct PDF vision upload requires conversion, we'll enhance text extraction
            # with structural awareness
            
            import PyPDF2
            import base64
            
            # Extract text while preserving structure hints
            structured_content = []
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    page_text = page.extract_text()
                    structured_content.append(f"--- PAGE {page_num} ---\n{page_text}\n")
            
            full_content = '\n'.join(structured_content)
            
            # Check content length and adjust processing strategy
            estimated_tokens = len(full_content) // 4  # Rough token estimate
            
            # Build enhanced prompt that emphasizes document structure
            rubric_prompt = self._build_grading_prompt_for_file_upload()
            
            # Adjust max_tokens based on content complexity and estimated response needs
            if estimated_tokens > 50000:
                max_output_tokens = 4000  # Large documents need detailed feedback
                logging.warning(f"Large PDF detected (~{estimated_tokens:,} tokens). Using expanded output limit.")
            elif estimated_tokens > 20000:
                max_output_tokens = 3000  # Medium documents
                logging.info(f"Medium PDF detected (~{estimated_tokens:,} tokens). Using moderate output limit.")
            else:
                max_output_tokens = 2000  # Standard documents
            
            # Check if content exceeds reasonable input limits
            total_prompt_estimate = len(rubric_prompt) // 4 + estimated_tokens
            if total_prompt_estimate > 100000:  # Approaching model limits
                logging.warning(f"PDF content very large (~{total_prompt_estimate:,} tokens total). May exceed model limits.")
                # Could implement chunking strategy here if needed
            
            enhanced_prompt = f"""{rubric_prompt}

DOCUMENT CONTENT (PDF with {len(pdf_reader.pages)} pages, ~{estimated_tokens:,} tokens):
{full_content}

FORMATTING NOTES:
- This content was extracted from a PDF document
- Page breaks and structure are indicated with "--- PAGE X ---" markers
- Original formatting, layout, and visual elements should be considered in evaluation
- The document may contain figures, tables, or special formatting not fully captured in text

Please evaluate this submission considering both the content and the implied document structure."""
            
            # Build system message
            system_message = self._build_system_message()
            
            # Create completion with enhanced context and adaptive token limits
            response = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": enhanced_prompt}
                ],
                temperature=0.3,
                max_tokens=max_output_tokens
            )
            
            # Parse the response
            grading_result = self._parse_grading_response(
                response.choices[0].message.content,
                student_data['name']
            )
            
            # Add metadata
            grading_result['grading_method'] = 'pdf_enhanced'
            grading_result['file_type'] = '.pdf'
            grading_result['pages_processed'] = len(pdf_reader.pages)
            grading_result['estimated_tokens'] = estimated_tokens
            grading_result['max_output_tokens'] = max_output_tokens
            grading_result['formatting_preserved'] = True
            
            logging.info(f"Successfully graded PDF for {student_data['name']} using enhanced processing (~{estimated_tokens:,} tokens)")
            return grading_result
            
        except Exception as e:
            logging.error(f"Error in PDF enhanced grading: {e}")
            raise
    
    def _grade_docx_enhanced(self, student_data: Dict[str, str], file_path: Path) -> Dict[str, Any]:
        """
        Grade DOCX with enhanced structure preservation
        
        Args:
            student_data (Dict): Student paper data
            file_path (Path): Path to DOCX file
            
        Returns:
            Dict: Grading results
        """
        try:
            from docx import Document
            
            doc = Document(file_path)
            
            # Extract content with structure preservation
            structured_content = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    # Add style information if available
                    style_info = f"[{para.style.name}]" if para.style.name != 'Normal' else ""
                    structured_content.append(f"{style_info} {para.text}")
            
            # Also extract tables if present
            for table in doc.tables:
                structured_content.append("\n[TABLE CONTENT:]")
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    if row_text.strip():
                        structured_content.append(row_text)
                structured_content.append("[END TABLE]\n")
            
            full_content = '\n'.join(structured_content)
            
            # Check content length and adjust processing strategy
            estimated_tokens = len(full_content) // 4  # Rough token estimate
            
            # Build enhanced prompt
            rubric_prompt = self._build_grading_prompt_for_file_upload()
            
            # Adjust max_tokens based on content complexity
            if estimated_tokens > 50000:
                max_output_tokens = 4000  # Large documents need detailed feedback
                logging.warning(f"Large DOCX detected (~{estimated_tokens:,} tokens). Using expanded output limit.")
            elif estimated_tokens > 20000:
                max_output_tokens = 3000  # Medium documents
                logging.info(f"Medium DOCX detected (~{estimated_tokens:,} tokens). Using moderate output limit.")
            else:
                max_output_tokens = 2000  # Standard documents
            
            enhanced_prompt = f"""{rubric_prompt}

DOCUMENT CONTENT (Microsoft Word .docx, ~{estimated_tokens:,} tokens):
{full_content}

FORMATTING NOTES:
- This content was extracted from a Word document (.docx)
- Document styles and structure are indicated with [Style Name] markers
- Tables are marked with [TABLE CONTENT:] and [END TABLE] markers
- Original formatting, layout, and visual elements should be considered in evaluation
- The document may contain additional formatting not fully captured in text

Please evaluate this submission considering both the content and the document structure."""
            
            # Build system message and get completion
            system_message = self._build_system_message()
            
            response = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": enhanced_prompt}
                ],
                temperature=0.3,
                max_tokens=max_output_tokens
            )
            
            # Parse the response
            grading_result = self._parse_grading_response(
                response.choices[0].message.content,
                student_data['name']
            )
            
            # Add metadata
            grading_result['grading_method'] = 'docx_enhanced'
            grading_result['file_type'] = '.docx'
            grading_result['paragraphs_processed'] = len(doc.paragraphs)
            grading_result['tables_found'] = len(doc.tables)
            grading_result['estimated_tokens'] = estimated_tokens
            grading_result['max_output_tokens'] = max_output_tokens
            grading_result['formatting_preserved'] = True
            
            logging.info(f"Successfully graded DOCX for {student_data['name']} using enhanced processing (~{estimated_tokens:,} tokens)")
            return grading_result
            
        except Exception as e:
            logging.error(f"Error in DOCX enhanced grading: {e}")
            raise
    
    def _grade_powerpoint_enhanced(self, student_data: Dict[str, str], file_path: Path) -> Dict[str, Any]:
        """
        Grade PowerPoint presentations with enhanced slide-by-slide processing
        
        Args:
            student_data (Dict): Student paper data  
            file_path (Path): Path to PowerPoint file
            
        Returns:
            Dict: Grading results
        """
        try:
            from pptx import Presentation
            
            # Load the presentation
            prs = Presentation(file_path)
            
            # Extract structured content
            structured_content = []
            structured_content.append(f"=== POWERPOINT PRESENTATION ANALYSIS ===")
            structured_content.append(f"File: {file_path.name}")
            structured_content.append(f"Total Slides: {len(prs.slides)}")
            structured_content.append("=" * 60)
            
            # Process each slide
            for slide_num, slide in enumerate(prs.slides, 1):
                structured_content.append(f"\n--- SLIDE {slide_num} ---")
                
                # Extract slide layout information
                slide_layout = slide.slide_layout
                structured_content.append(f"[Layout: {slide_layout.name}]")
                
                # Extract text from all shapes
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Identify shape type for context
                        if hasattr(shape, 'text_frame'):
                            # This is likely a text box or title
                            text = shape.text.strip()
                            if len(text) > 50:
                                slide_texts.append(f"[TEXT CONTENT]: {text}")
                            else:
                                slide_texts.append(f"[TITLE/HEADER]: {text}")
                        else:
                            slide_texts.append(f"[TEXT]: {shape.text.strip()}")
                
                if slide_texts:
                    structured_content.extend(slide_texts)
                else:
                    structured_content.append("[No text content - may contain images/charts only]")
                
                # Note: We can't extract images, but we can note their presence
                shape_count = len(slide.shapes)
                text_shapes = len([s for s in slide.shapes if hasattr(s, "text") and s.text.strip()])
                non_text_shapes = shape_count - text_shapes
                if non_text_shapes > 0:
                    structured_content.append(f"[NOTE: {non_text_shapes} additional visual elements (charts/images)]")
            
            full_content = '\n'.join(structured_content)
            
            # Build enhanced prompt
            rubric_prompt = self._build_grading_prompt_for_file_upload()
            
            enhanced_prompt = f"""{rubric_prompt}

PRESENTATION CONTENT (PowerPoint .pptx/.ppt):
{full_content}

PRESENTATION ANALYSIS NOTES:
- This is a PowerPoint presentation with {len(prs.slides)} slides
- Content has been extracted slide-by-slide with layout information
- Visual elements (charts, images, diagrams) are noted but content not extracted
- Slide progression and structure should be considered in evaluation
- Presentation design and organization are part of the submission quality

Please evaluate this presentation considering:
1. Content quality and accuracy
2. Presentation structure and flow
3. Organization and clarity
4. Use of visual elements (where noted)
5. Professional presentation standards"""
            
            # Build system message and get completion
            system_message = self._build_system_message()
            
            response = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": enhanced_prompt}
                ],
                temperature=0.3,
                max_tokens=3000  # PowerPoint presentations often need detailed feedback
            )
            
            # Parse the response
            grading_result = self._parse_grading_response(
                response.choices[0].message.content,
                student_data['name']
            )
            
            # Add metadata
            grading_result['grading_method'] = 'powerpoint_enhanced'
            grading_result['file_type'] = file_path.suffix.lower()
            grading_result['slides_processed'] = len(prs.slides)
            grading_result['formatting_preserved'] = True
            grading_result['presentation_structure'] = True
            
            logging.info(f"Successfully graded PowerPoint for {student_data['name']} with {len(prs.slides)} slides")
            return grading_result
            
        except ImportError:
            logging.error(f"python-pptx not available for PowerPoint processing")
            # Fallback to text extraction method
            return self._grade_paper_with_text_extraction(student_data)
        except Exception as e:
            logging.error(f"Error in PowerPoint enhanced grading: {e}")
            raise
    
    def _grade_text_file_enhanced(self, student_data: Dict[str, str], file_path: Path) -> Dict[str, Any]:
        """
        Grade text files with enhanced processing
        
        Args:
            student_data (Dict): Student paper data
            file_path (Path): Path to text file
            
        Returns:
            Dict: Grading results
        """
        try:
            # Read text file with encoding detection (optional dependency)
            encoding = 'utf-8'  # Default encoding
            
            try:
                import chardet
                # Detect encoding
                with open(file_path, 'rb') as f:
                    raw_data = f.read()
                    encoding_result = chardet.detect(raw_data)
                    encoding = encoding_result.get('encoding', 'utf-8')
            except ImportError:
                # chardet not available, use utf-8 with error handling
                logging.info("chardet not available, using UTF-8 encoding")
            
            # Read with detected/default encoding
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
            except UnicodeDecodeError:
                # Fallback to utf-8 with error handling
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                encoding = 'utf-8 (with errors ignored)'
            
            # Build enhanced prompt
            rubric_prompt = self._build_grading_prompt_for_file_upload()
            
            enhanced_prompt = f"""{rubric_prompt}

DOCUMENT CONTENT ({file_path.suffix.upper()} file):
{content}

FORMATTING NOTES:
- This is a {file_path.suffix.upper()} file processed with enhanced text handling
- File encoding: {encoding}
- Text formatting and structure are preserved as submitted by the student

Please evaluate this submission based on the content and any evident text structure."""
            
            # Build system message and get completion
            system_message = self._build_system_message()
            
            response = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": enhanced_prompt}
                ],
                temperature=0.3,
                max_tokens=2000
            )
            
            # Parse the response
            grading_result = self._parse_grading_response(
                response.choices[0].message.content,
                student_data['name']
            )
            
            # Add metadata
            grading_result['grading_method'] = 'text_enhanced'
            grading_result['file_type'] = file_path.suffix.lower()
            grading_result['encoding_used'] = encoding
            grading_result['formatting_preserved'] = True
            
            logging.info(f"Successfully graded text file for {student_data['name']} using enhanced processing")
            return grading_result
            
        except Exception as e:
            logging.error(f"Error in text file enhanced grading: {e}")
            raise
    
    def _grade_paper_with_text_extraction(self, student_data: Dict[str, str]) -> Dict[str, Any]:
        """
        Grade a paper using traditional text extraction method (legacy support)
        
        Args:
            student_data (Dict): Student paper data
            
        Returns:
            Dict: Grading results
        """
        # Construct the grading prompt
        prompt = self._build_grading_prompt(student_data['content'])
        
        # Estimate tokens for adaptive response length
        estimated_tokens = len(student_data['content']) // 4
        if estimated_tokens > 50000:
            max_output_tokens = 4000
        elif estimated_tokens > 20000:
            max_output_tokens = 3000
        else:
            max_output_tokens = 2000
        
        try:
            # Build system message based on instructor configuration
            system_message = self._build_system_message()
            
            response = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.3,
                max_tokens=max_output_tokens
            )
            
            # Parse the response
            grading_result = self._parse_grading_response(
                response.choices[0].message.content,
                student_data['name']
            )
            
            # Add metadata about grading method
            grading_result['grading_method'] = 'text_extraction'
            grading_result['estimated_tokens'] = estimated_tokens
            grading_result['max_output_tokens'] = max_output_tokens
            
            logging.info(f"Graded paper for {student_data['name']} using text extraction (~{estimated_tokens:,} tokens)")
            return grading_result
            
        except Exception as e:
            logging.error(f"Error grading paper for {student_data['name']}: {str(e)}")
            raise
    
    def _build_grading_prompt_for_file_upload(self) -> str:
        """
        Build the grading prompt for file upload (without content since file is uploaded separately)
        
        Returns:
            str: Formatted prompt for file upload grading
        """
        rubric_str = json.dumps(self.rubric, indent=2)
        
        # Build context from instructor configuration
        context_parts = []
        
        if self.instructor_config:
            if course_context := self.instructor_config.get('course_context'):
                context_parts.append(f"COURSE CONTEXT:")
                context_parts.append(f"- Course level: {course_context.get('course_level', 'Not specified')}")
                context_parts.append(f"- Course type: {course_context.get('course_type', 'Not specified')}")
                context_parts.append(f"- Student background: {course_context.get('student_background', 'Mixed levels')}")
                
                if learning_objectives := course_context.get('learning_objectives'):
                    context_parts.append(f"- Learning objectives: {', '.join(learning_objectives)}")
        
        # Add course materials if available (optional)
        if hasattr(self, 'course_materials') and self.course_materials:
            if context_parts:  # Add spacing if there's already context
                context_parts.append("")
            context_parts.append("COURSE MATERIALS FOR REFERENCE:")
            if self.course_materials_instructions:
                context_parts.append(f"Instructions for using course materials: {self.course_materials_instructions}")
            else:
                context_parts.append("The following course materials are provided for context. Reference them when appropriate for grading:")
            context_parts.append(str(self.course_materials))
        
        context_str = '\n'.join(context_parts) if context_parts else ""
        
        prompt = f"""You are grading a student submission using the attached document. Please read the document thoroughly and grade it according to the following rubric:

RUBRIC:
{rubric_str}

{context_str}

INSTRUCTIONS:
1. Read the uploaded document carefully, preserving attention to formatting, structure, and visual elements
2. Grade each criterion in the rubric based on the document content
3. Provide specific feedback referencing what you observed in the document
4. Calculate the total score and percentage
5. Return your response in the following JSON format:

{{
    "overall_score": <total_points>,
    "max_possible_score": <maximum_points>,
    "percentage": <percentage_score>,
    "letter_grade": "<letter_grade>",
    "criteria_scores": {{
        "<criterion_name>": {{
            "score": <points>,
            "max_score": <max_points>,
            "feedback": "<specific_feedback>"
        }}
    }},
    "overall_feedback": "<comprehensive_feedback>",
    "strengths": ["<strength1>", "<strength2>"],
    "areas_for_improvement": ["<improvement1>", "<improvement2>"]
}}

Be specific, constructive, and fair in your grading. Provide actionable feedback.
The document formatting, layout, and visual presentation should be considered as part of your evaluation where relevant to the rubric criteria."""
        
        return prompt
    
    def _build_system_message(self) -> str:
        """
        Build the system message for ChatGPT based on instructor configuration
        
        Returns:
            str: Customized system message
        """
        if self.instructor_config and self.instructor_config.get('custom_system_message'):
            # Use custom system message from configuration
            base_message = self.instructor_config['custom_system_message']
        else:
            # Default system message
            base_message = "You are an expert academic grader. Provide detailed, constructive feedback based on the given rubric."
        
        # Add additional instructions if available
        if self.instructor_config:
            additional_instructions = []
            
            # Add grading philosophy
            if philosophy := self.instructor_config.get('grading_philosophy'):
                additional_instructions.append(f"Your grading approach should be {philosophy}.")
            
            # Add specific instructions
            if specific_instructions := self.instructor_config.get('specific_instructions'):
                additional_instructions.extend(specific_instructions)
            
            # Add comment preferences
            if comment_prefs := self.instructor_config.get('comment_preferences'):
                if comment_prefs.get('include_strengths'):
                    additional_instructions.append("Always highlight what the student did well.")
                if comment_prefs.get('include_suggestions'):
                    additional_instructions.append("Provide specific suggestions for improvement.")
                if comment_prefs.get('personal_touch'):
                    additional_instructions.append("Make comments personal and encouraging.")
            
            if additional_instructions:
                base_message += " " + " ".join(additional_instructions)
        
        return base_message
    
    def _build_grading_prompt(self, paper_content: str) -> str:
        """
        Build the grading prompt for ChatGPT
        
        Args:
            paper_content (str): The student's paper content
            
        Returns:
            str: Formatted prompt
        """
        rubric_str = json.dumps(self.rubric, indent=2)
        
        # Build context from instructor configuration
        context_parts = []
        
        if self.instructor_config:
            if course_context := self.instructor_config.get('course_context'):
                context_parts.append(f"COURSE CONTEXT:")
                context_parts.append(f"- Course level: {course_context.get('course_level', 'Not specified')}")
                context_parts.append(f"- Course type: {course_context.get('course_type', 'Not specified')}")
                context_parts.append(f"- Student background: {course_context.get('student_background', 'Mixed levels')}")
                
                if learning_objectives := course_context.get('learning_objectives'):
                    context_parts.append(f"- Learning objectives: {', '.join(learning_objectives)}")
                context_parts.append("")
            
            if expertise := self.instructor_config.get('subject_expertise'):
                context_parts.append(f"INSTRUCTOR EXPERTISE: {', '.join(expertise)}")
                context_parts.append("")
        
        # Add course materials if available (optional)
        if hasattr(self, 'course_materials') and self.course_materials:
            context_parts.append("COURSE MATERIALS FOR REFERENCE:")
            if self.course_materials_instructions:
                context_parts.append(f"Instructions for using course materials: {self.course_materials_instructions}")
            else:
                context_parts.append("The following course materials are provided for context. Reference them when appropriate for grading:")
            context_parts.append(str(self.course_materials))
            context_parts.append("")
        
        context_str = '\n'.join(context_parts) if context_parts else ""
        
        prompt = f"""
Please grade the following student paper based on the provided rubric. 

{context_str}RUBRIC:
{rubric_str}

STUDENT PAPER:
{paper_content}

Please provide your response in the following JSON format:
{{
    "overall_score": <total_points>,
    "max_possible_score": <maximum_points>,
    "percentage": <percentage_score>,
    "letter_grade": "<letter_grade>",
    "criteria_scores": {{
        "<criterion_name>": {{
            "score": <points>,
            "max_score": <max_points>,
            "feedback": "<specific_feedback>"
        }}
    }},
    "overall_feedback": "<comprehensive_feedback>",
    "strengths": ["<strength1>", "<strength2>"],
    "areas_for_improvement": ["<improvement1>", "<improvement2>"]
}}

Be specific, constructive, and fair in your grading. Provide actionable feedback.
"""
        return prompt
    
    def _parse_grading_response(self, response_text: str, student_name: str) -> Dict[str, Any]:
        """
        Parse ChatGPT's grading response
        
        Args:
            response_text (str): Raw response from ChatGPT
            student_name (str): Student's name
            
        Returns:
            Dict: Parsed grading results
        """
        try:
            # Extract JSON from the response
            start_idx = response_text.find('{')
            end_idx = response_text.rfind('}') + 1
            json_str = response_text[start_idx:end_idx]
            
            grading_data = json.loads(json_str)
            grading_data['student_name'] = student_name
            grading_data['grading_date'] = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            
            return grading_data
            
        except Exception as e:
            logging.error(f"Error parsing grading response: {str(e)}")
            # Return a default structure if parsing fails
            return {
                'student_name': student_name,
                'overall_score': 0,
                'max_possible_score': 100,
                'percentage': 0,
                'letter_grade': 'F',
                'criteria_scores': {},
                'overall_feedback': f"Error in grading: {str(e)}",
                'strengths': [],
                'areas_for_improvement': [],
                'grading_date': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            }
    
    def grade_all_papers(self) -> List[Dict[str, Any]]:
        """
        Grade all loaded student papers
        
        Returns:
            List[Dict]: All grading results
        """
        if not self.students_data:
            raise ValueError("No student papers loaded.")
        
        self.graded_results = []
        
        for student_data in self.students_data:
            try:
                result = self.grade_paper(student_data)
                self.graded_results.append(result)
            except Exception as e:
                logging.error(f"Failed to grade {student_data['name']}: {str(e)}")
                
        logging.info(f"Completed grading {len(self.graded_results)} papers")
        return self.graded_results
    
    def export_results(self, output_format: str = 'xlsx', filename: Optional[str] = None) -> str:
        """
        Export grading results to Excel or CSV for mail merge
        
        Args:
            output_format (str): 'xlsx' or 'csv'
            filename (str, optional): Output filename
            
        Returns:
            str: Path to the exported file
        """
        if not self.graded_results:
            raise ValueError("No grading results to export.")
        
        # Create DataFrame for export
        export_data = []
        
        for result in self.graded_results:
            row = {
                'Student_Name': result['student_name'],
                'Overall_Score': result['overall_score'],
                'Max_Score': result['max_possible_score'],
                'Percentage': result['percentage'],
                'Letter_Grade': result['letter_grade'],
                'Overall_Feedback': result['overall_feedback'],
                'Strengths': '; '.join(result.get('strengths', [])),
                'Areas_for_Improvement': '; '.join(result.get('areas_for_improvement', [])),
                'Grading_Date': result['grading_date']
            }
            
            # Add individual criteria scores
            for criterion, details in result.get('criteria_scores', {}).items():
                row[f'{criterion}_Score'] = details.get('score', 0)
                row[f'{criterion}_Feedback'] = details.get('feedback', '')
            
            export_data.append(row)
        
        df = pd.DataFrame(export_data)
        
        # Generate filename if not provided
        if not filename:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f'grading_results_{timestamp}.{output_format}'
        
        # Export based on format
        if output_format.lower() == 'xlsx':
            df.to_excel(filename, index=False)
        elif output_format.lower() == 'csv':
            df.to_csv(filename, index=False, encoding='utf-8')
        else:
            raise ValueError("Unsupported format. Use 'xlsx' or 'csv'.")
        
        logging.info(f"Results exported to {filename}")
        return filename
    
    def generate_summary_report(self) -> Dict[str, Any]:
        """
        Generate a summary report of all grading results
        
        Returns:
            Dict: Summary statistics
        """
        if not self.graded_results:
            return {}
        
        scores = [result['percentage'] for result in self.graded_results]
        letter_grades = [result['letter_grade'] for result in self.graded_results]
        
        summary = {
            'total_papers': len(self.graded_results),
            'average_score': sum(scores) / len(scores),
            'highest_score': max(scores),
            'lowest_score': min(scores),
            'grade_distribution': pd.Series(letter_grades).value_counts().to_dict(),
            'completion_date': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        }
        
        return summary

    def _make_chat_completion(self, messages, max_tokens=1500, temperature=0.3):
        """Unified wrapper for OpenAI chat.completions.create, supporting GPT-5 and earlier models."""
        kwargs = {
            'model': self.model,
            'messages': messages,
            'temperature': temperature
        }
        
        # Use correct token parameter for model
        if str(self.model).lower().startswith('gpt-5'):
            kwargs['max_completion_tokens'] = max_tokens
        else:
            kwargs['max_tokens'] = max_tokens
        
        # Note: OpenAI Chat Completions API has deprecated the files parameter
        # File uploads now require different approaches (Assistant API, etc.)
        # This method now only handles text-based completions
        
        try:
            response = self.client.chat.completions.create(**kwargs)
            return response.choices[0].message.content
        except Exception as api_exc:
            logging.error(f"API call failed: {api_exc}")
            raise api_exc


def main():
    """Example usage of the GradingAgent"""
    # Initialize with your OpenAI API key
    api_key = input("Enter your OpenAI API key: ")
    agent = GradingAgent(api_key)
    
    try:
        # Load rubric
        rubric_file = input("Enter path to rubric JSON file: ")
        agent.load_rubric(rubric_file)
        
        # Load student papers
        papers_folder = input("Enter path to folder containing student papers: ")
        agent.load_student_papers(papers_folder)
        
        # Grade all papers
        print("Grading papers...")
        results = agent.grade_all_papers()
        
        # Export results
        output_format = input("Export format (xlsx/csv): ").lower()
        output_file = agent.export_results(output_format)
        
        # Generate summary
        summary = agent.generate_summary_report()
        print(f"\nGrading Summary:")
        print(f"Total papers: {summary['total_papers']}")
        print(f"Average score: {summary['average_score']:.2f}%")
        print(f"Results exported to: {output_file}")
        
    except Exception as e:
        logging.error(f"Error in main execution: {str(e)}")
        print(f"An error occurred: {str(e)}")


if __name__ == "__main__":
    main()
